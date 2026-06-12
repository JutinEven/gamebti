"""
文档读取工具 (Read Document Tool)

使用本地解析库提取文档内容，支持 PDF/Word/Excel/PPT/TXT/Markdown。
替代原 Coze Plugin API 依赖。

依赖库（已安装）：
- pypdf:       PDF 解析
- docx2python: Word (.docx) 解析
- openpyxl:    Excel (.xlsx) 解析
- python-pptx: PowerPoint (.pptx) 解析
- chardet:     文本编码检测
"""

import os
import json
import logging
from langchain.tools import tool

logger = logging.getLogger(__name__)

# 支持的文件扩展名
SUPPORTED_EXTS = {
    ".pdf": "PDF 文档",
    ".docx": "Word 文档",
    ".doc": "Word 文档 (旧格式)",
    ".xlsx": "Excel 表格",
    ".xls": "Excel 表格 (旧格式)",
    ".csv": "CSV 表格",
    ".pptx": "PowerPoint 演示",
    ".ppt": "PowerPoint 演示 (旧格式)",
    ".txt": "纯文本",
    ".md": "Markdown",
    ".json": "JSON 数据",
    ".html": "HTML 网页",
    ".htm": "HTML 网页",
}


@tool
def read_document(file_path: str) -> str:
    """
    读取本地文档内容，支持 PDF、Word、Excel、PowerPoint、文本等格式。

    当用户上传或引用文档时，使用此工具提取文档中的文字内容。

    Args:
        file_path: 文档的本地文件路径（绝对路径或相对于工作目录的路径）

    Returns:
        文档的文本内容
    """
    if not file_path:
        return "错误：未提供文件路径"

    # 检查文件是否存在
    if not os.path.exists(file_path):
        # 尝试在工作目录下查找
        alt_path = os.path.join(os.getcwd(), file_path)
        if os.path.exists(alt_path):
            file_path = alt_path
        else:
            return f"错误：文件不存在 — {file_path}"

    # 获取扩展名
    _, ext = os.path.splitext(file_path)
    ext = ext.lower()

    if ext not in SUPPORTED_EXTS:
        return (
            f"不支持的文件格式: {ext}\n"
            f"支持的格式: {', '.join(SUPPORTED_EXTS.keys())}"
        )

    logger.info(f"读取文档: {file_path} ({SUPPORTED_EXTS.get(ext, ext)})")

    try:
        content = _extract_content(file_path, ext)
        return content if content else "文档内容为空。"
    except ImportError as e:
        return f"解析库缺失: {e}。请安装对应的解析库。"
    except Exception as e:
        logger.error(f"文档解析失败: {e}")
        return f"文档解析失败: {e}"


def _extract_content(file_path: str, ext: str) -> str:
    """根据文件类型调用对应的解析器"""
    if ext == ".pdf":
        return _read_pdf(file_path)
    elif ext in (".docx", ".doc"):
        return _read_docx(file_path)
    elif ext in (".xlsx", ".xls", ".csv"):
        return _read_spreadsheet(file_path, ext)
    elif ext in (".pptx", ".ppt"):
        return _read_pptx(file_path)
    elif ext in (".txt", ".md", ".json", ".html", ".htm"):
        return _read_text(file_path)
    else:
        return f"不支持的格式: {ext}"


def _read_pdf(file_path: str) -> str:
    """读取 PDF 文件"""
    import pypdf

    reader = pypdf.PdfReader(file_path)
    pages_text = []
    for i, page in enumerate(reader.pages):
        text = page.extract_text()
        if text and text.strip():
            pages_text.append(f"--- 第 {i + 1} 页 ---\n{text.strip()}")
    return "\n\n".join(pages_text)


def _read_docx(file_path: str) -> str:
    """读取 Word 文档"""
    from docx2python import docx2python

    doc = docx2python(file_path)
    all_parts = []

    for section in doc.body:
        if isinstance(section, list):
            for item in section:
                if isinstance(item, list):
                    for sub_item in item:
                        if isinstance(sub_item, str) and sub_item.strip():
                            all_parts.append(sub_item.strip())
                        elif isinstance(sub_item, list):
                            row_text = "\n".join(
                                str(cell).strip()
                                for cell in sub_item
                                if str(cell).strip()
                            )
                            if row_text:
                                all_parts.append(row_text)
                elif isinstance(item, str) and item.strip():
                    all_parts.append(item.strip())

    doc.close()
    return "\n\n".join(all_parts)


def _read_spreadsheet(file_path: str, ext: str) -> str:
    """读取 Excel/CSV 文件"""
    import pandas as pd

    if ext == ".csv":
        df = pd.read_csv(file_path)
    else:
        df = pd.read_excel(file_path)

    return df.to_string()


def _read_pptx(file_path: str) -> str:
    """读取 PowerPoint 文件"""
    from pptx import Presentation

    prs = Presentation(file_path)
    full_text = []

    for i, slide in enumerate(prs.slides):
        page_content = [f"=== 第 {i + 1} 页 ==="]

        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                page_content.append(shape.text.strip())

            if shape.has_table:
                table_texts = []
                for row in shape.table.rows:
                    row_cells = [
                        cell.text_frame.text.strip()
                        for cell in row.cells
                        if cell.text_frame.text.strip()
                    ]
                    if row_cells:
                        table_texts.append(" | ".join(row_cells))
                if table_texts:
                    page_content.append("[表格]\n" + "\n".join(table_texts))

        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text
            if notes.strip():
                page_content.append(f"[备注]: {notes.strip()}")

        full_text.append("\n".join(page_content))

    return "\n\n".join(full_text)


def _read_text(file_path: str) -> str:
    """读取纯文本文件（自动检测编码）"""
    import chardet

    with open(file_path, "rb") as f:
        raw = f.read()

    detected = chardet.detect(raw)
    encoding = detected.get("encoding", "utf-8") if detected else "utf-8"

    return raw.decode(encoding, errors="replace")
